from pptx import Presentation

apresentacao = Presentation()

# editar o ppt
slide1 = apresentacao.slides.add_slide(apresentacao.slide_layouts[0])

titulo = slide1.shapes.title
subtitulo = slide1.placeholders[1]

titulo.text = "1º Slide do Gui"
subtitulo.text = "Estamos criando ppt com Python"

#salvar o ppt
apresentacao.save(f"P:\Guilherme\Python\MeuPPT.pptx")

# Criar personalizados

from pptx import Presentation
from pptx.util import Inches

apresentacao = Presentation()

x = Inches(1)
y = Inches(1)
largura = Inches(2)
altura = Inches(2)

slide = apresentacao.slides.add_slide(apresentacao.slide_layouts[6])

caixa_texto = slide.shapes.add_textbox(x, y, largura, altura)

caixa_texto.text = "Slide do Gui do Python e do Excel"

text_frame = caixa_texto.text_frame
paragrafo = text_frame.add_paragraph()
paragrafo.text = "Gui do Excel e do Python é um cara que busca evoluir no TI"
paragrafo.font.bold = True


apresentacao.save(f"P:\Guilherme\Python\MeuPPT2.pptx")

# Criar com gráfico

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

apresentacao = Presentation()

slide = apresentacao.slides.add_slide(apresentacao.slide_layouts[6])

produtos = ["Iphone", "Ipad", "Airpod"]
vendas = [1500, 1000, 2000]

x = Inches(1)
y = Inches(1)
largura = Inches(7)
altura = Inches(6)

dados_grafico = CategoryChartData()
dados_grafico.categories = produtos
dados_grafico.add_series("Vendas", vendas)
slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, largura, altura, dados_grafico)

apresentacao.save(f"P:\Guilherme\Python\MeuPPT3.pptx")